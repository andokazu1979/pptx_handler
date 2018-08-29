#!/usr/bin/env python
# -*- coding: utf-8 -*-

import datetime
import re
import sys
import logging

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from toml_parser import TOMLParser

parser = TOMLParser()
args = sys.argv
if len(args) == 1:
    raise Exception("Specify configuration file!")
elif len(args) >= 3:
    raise Exception("Too many configuration files is specified!")
parser.parse(sys.argv[1])
conf = parser.dict_root
project = conf['global']['project']
loglevel = conf['global']['loglevel']

if loglevel == 'DEBUG':
    level_ = logging.DEBUG
elif loglevel == 'INFO':
    level_ = logging.INFO
elif loglevel == 'WARNING':
    level_ = logging.WARNING
elif loglevel == 'ERROR':
    level_ = logging.ERROR
elif loglevel == 'CRITCAL':
    level_ = logging.CRITCAL

logging.basicConfig(level = level_)
logger = logging.getLogger(__name__)
logger.info(project)
project_conf = conf[project]
logger.debug(project_conf)

class PPTXHandler(object):
    def __init__(self):

        self.prs = Presentation()

        self.dirpath_in = project_conf['dirpath_in']
        self.dirpath_out = project_conf['dirpath_out']
        self.delta_t = project_conf['delta_t']

        # Parameters
        self.ncols = project_conf['ncols']
        self.sizex = project_conf['sizex']
        self.sizey = project_conf['sizey']
        self.col_int = project_conf['col_int']
        self.row_int = project_conf['row_int'] 
        self.col_sta = project_conf['col_sta'] 
        self.row_sta = project_conf['row_sta'] 

        self.left = project_conf['left']
        self.top = project_conf['top']
        self.width = project_conf['width']
        self.height = project_conf['height']

        self.title = project_conf['title']
        self.str_suffix = project_conf['str_suffix']

        self.lst_slide = project_conf['slide']

        self.lst_fig_category = project_conf['lst_fig_category']

    def create_pptx(self):
        self.const_description_slides()
        self.const_data_slides()
        self.output_pptx()

    def output_pptx(self):
        dt_now = datetime.datetime.now()
        self.prs.save(u'{4}/{0:04d}{1:02d}{2:02d}_{3}.pptx'.format(
            dt_now.year, dt_now.month, dt_now.day, self.title, self.dirpath_out))

    def const_description_slides(self):
        for conf_slide in self.lst_slide:
            self.const_slide(conf_slide)

    def get_lst_dt(self, dt1, dt2, delta_t):
        period = dt2 - dt1
        days = period.days
        hours = period.seconds / 3600
        ngrids_time = (days * 24 + hours) / delta_t + 1
        td = datetime.timedelta(hours=delta_t)
        return [dt1 + td * i for i in range(ngrids_time)]

    def const_data_slides(self):
        for fig_category in self.lst_fig_category:
            self.fig_category = fig_category
            self.category_conf = project_conf['exec_cond_{1}'.format(project, fig_category)]
            lst_period = []
            for period_sta, period_end in zip(self.category_conf['lst_period_sta'], self.category_conf['lst_period_end']):
                lst_period.append([datetime.datetime.strptime(period_sta, '%Y%m%d%H'), datetime.datetime.strptime(period_end, '%Y%m%d%H')])
            for i, period in enumerate(lst_period):
                self.target = self.category_conf['targets'][i]
                self.care_for_each_period(period)
                for dt in self.get_lst_dt(period[0], period[1], self.delta_t):
                    print(dt)
                    self.loop(dt)

    def do_inner_proc(self, dt, slide):
        lst_fig_title = self.category_conf['lst_fig_title']
        for i, fig_type in enumerate(self.category_conf['lst_fig_type']):
            fig_dir = self.category_conf['lst_fig_dir'][i]
            icol = i % self.ncols
            irow = i / self.ncols

            if self.special_care(fig_type, icol, irow, dt, slide):
                continue

            if fig_type == "-":
                continue

            left = Inches(self.col_sta + self.col_int * icol)
            top = Inches(self.row_sta + self.row_int * irow)
            width = Inches(self.sizex)
            height = Inches(self.sizey)
            if self.str_suffix == '{0:04d}{1:02d}{2:02d}{3:02d}.png':
                self.suffix = self.str_suffix.format(dt.year, dt.month, dt.day, dt.hour)
            elif self.str_suffix == '{3:02d}Z{2:02d}{1}{0:04d}.png':
                self.suffix = self.str_suffix.format(dt.year, dt.strftime("%b").upper(), dt.day, dt.hour)
            img_path = self.get_filepath(fig_dir, fig_type, self.suffix)
            slide.shapes.add_picture(img_path, left, top, height, width)

            top_figtitle = top - Inches(0.020)
            height_figtitle = Inches(0.4)
            txBox_figtitle = slide.shapes.add_textbox(left, top_figtitle, width, height_figtitle)
            txBox_figtitle.fill.solid()
            txBox_figtitle.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf_figtitle = txBox_figtitle.text_frame
            tf_figtitle.text = lst_fig_title[i]

            left_ = Inches(9.45-len(self.section_title)*0.09225)
            top_ = Inches(0.1)
            width_ = Inches(0.45+len(self.section_title)*0.0925)
            height_ = Inches(0.4)
            txBox = slide.shapes.add_textbox(left_, top_, width_, height_)
            txBox.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
            tf = txBox.text_frame
            tf.text = self.section_title

    def const_slide(self, conf_slide):
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        shapes.title.text = conf_slide['title']

        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = conf_slide['text1']
        self.bullet(tf, conf_slide, 'text2', 'level2')
        self.bullet(tf, conf_slide, 'text3', 'level3')
        self.bullet(tf, conf_slide, 'text4', 'level4')
        self.bullet(tf, conf_slide, 'text5', 'level5')
        self.bullet(tf, conf_slide, 'text6', 'level6')
        self.bullet(tf, conf_slide, 'text7', 'level7')

    def bullet(self, tf, conf_slide, key_text, key_level):
        if key_text in conf_slide:
            p = tf.add_paragraph()
            p.text = conf_slide[key_text]
            p.level = conf_slide[key_level]

    def get_filepath(self, fig_dir, fig_type, suffix):
        if fig_type != '-':
            img_path = '{0}/{1}/{2}/{3}/{2}_{4}'.format(self.dirpath_in, fig_dir, fig_type, self.target, suffix)
        return img_path

    def loop(self, dt):
        title_only_slide_layout = self.prs.slide_layouts[5]

        slide = self.prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes
        shapes.title.text = dt.strftime('%Y/%m/%d %HUTC')
        self.do_inner_proc(dt, slide)

    def care_for_each_period(self, period):
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        self.section_title = '{0} {1}'.format(self.fig_category, self.target)
        shapes.title.text = self.section_title
        days = (period[1] - period[0]).days
        hours = (period[1] - period[0]).seconds / 3600
        if days == 0:
            str_days = "{}hours".format(hours)
        else:
            str_days = "{}days".format(days)
        slide.placeholders[1].text = '{0} - {1} ({2})'.format(period[0].strftime(
            "%Y/%m/%d %HUTC"), period[1].strftime("%Y/%m/%d %HUTC"), str_days)

    def special_care(self, fig_type, icol, irow, dt, slide):
        if fig_type[-1] == '_':
            left = Inches(self.left)
            top = Inches(self.top)
            width = Inches(self.width)
            height = Inches(self.height)
            img_path = '{0}/{1}/{2}/{1}_{3}'.format(self.dirpath_in, fig_type[:-1], self.target, self.suffix)
            print(img_path)
            try:
                slide.shapes.add_picture(img_path, left, top, height, width)
            except IOError as e:
                pass
            return True
        else:
            return False
        return False

    def const_second_slide(self):
        pass


if __name__ == '__main__':
    obj = PPTXHandler()
    obj.create_pptx()
